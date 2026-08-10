from __future__ import annotations

import json
import zipfile
from pathlib import Path
from typing import Any

from app.models.content_type import ContentType
from app.models.extracted_content import ExtractedContent
from app.models.attachment_processing_result import AttachmentProcessingStatus
from app.utils.logger import logger


class AttachmentContentExtractor:
    """Extract searchable text from downloaded Confluence attachments.

    Extraction is deliberately conservative for binary data, while unknown
    extensions are treated as text when the payload looks textual. JSON files
    also have a raw-text fallback because the Confluence repository contains
    many JSON-like configuration files that are not strict JSON.
    """

    MAX_ARCHIVE_FILES = 10_000
    MAX_ARCHIVE_BYTES = 512 * 1024 * 1024
    TEXT_SAMPLE_BYTES = 64 * 1024

    def __init__(self):
        self.download_folder = Path("downloads")
        self.extractors = self._build_extractor_registry()

    def _build_extractor_registry(self):
        return {
            ".xlsx": self._extract_excel,
            ".csv": self._extract_csv,
            ".json": self._extract_json,
            ".zip": self._extract_zip,
            ".rar": self._extract_rar,
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".txt": self._extract_text,
            ".xml": self._extract_text,
            ".html": self._extract_text,
            ".htm": self._extract_text,
            ".ldif": self._extract_ldif,
            ".png": self._extract_image,
            ".jpg": self._extract_image,
            ".jpeg": self._extract_image,
            ".gif": self._extract_image,
            ".bmp": self._extract_image,
            ".tif": self._extract_image,
            ".tiff": self._extract_image,
        }

    def supports(self, filename: str) -> bool:
        """Return whether the file can be handled by a deterministic extractor."""
        name = str(filename or "")
        suffix = Path(name).suffix.lower()
        if suffix in self.extractors:
            return True
        lowered = name.lower()
        return ".ldif-" in lowered or self._looks_like_filename_for_text(name)

    def extract_with_status(
        self,
        attachment: Any,
        file_path: str | Path | None = None,
    ) -> tuple[str, ExtractedContent | None, str | None]:
        filename = str(getattr(attachment, "filename", "<unknown>"))
        path = self._resolve_path(attachment, file_path)

        if not path.exists():
            error = f"File not found: {path}"
            logger.warning(error)
            return AttachmentProcessingStatus.EXTRACTION_FAILED, None, error

        try:
            extractor = self._resolve_extractor(filename, path)
            if extractor is None:
                error = f"No extractor for {Path(filename).suffix.lower() or '<no extension>'}"
                logger.warning(error)
                return AttachmentProcessingStatus.UNSUPPORTED, None, error

            extracted = extractor(path)
            if extracted is None:
                error = "Extractor returned no content"
                return AttachmentProcessingStatus.EXTRACTION_FAILED, None, error

            return AttachmentProcessingStatus.SUCCESS, extracted, None
        except Exception as exc:
            logger.exception(f"Failed extracting {filename}")
            return AttachmentProcessingStatus.EXTRACTION_FAILED, None, str(exc)

    def extract(self, attachment, file_path: str | Path | None = None):
        """Backward-compatible extraction API returning only content/None."""
        status, extracted, _ = self.extract_with_status(attachment, file_path)
        return extracted if status == AttachmentProcessingStatus.SUCCESS else None

    def _resolve_path(self, attachment: Any, file_path: str | Path | None) -> Path:
        if file_path is not None:
            return Path(file_path)
        return self.download_folder / str(getattr(attachment, "filename", ""))

    def _resolve_extractor(self, filename: str, file_path: Path):
        suffix = file_path.suffix.lower()
        if suffix in self.extractors:
            return self.extractors[suffix]

        lowered = filename.lower()
        if ".ldif-" in lowered:
            return self._extract_ldif

        if suffix == "":
            # Some historical Confluence attachment names contain no extension.
            # Detect common JSON first, then fall back to text if the payload is textual.
            if self._looks_like_text(file_path):
                return self._extract_unknown_text

        return None

    @staticmethod
    def _looks_like_filename_for_text(name: str) -> bool:
        # Names with no extension are handled by content sniffing; this method
        # exists to keep supports() useful without reading the file.
        return Path(name).suffix == ""

    def _looks_like_text(self, file_path: Path) -> bool:
        sample = file_path.read_bytes()[: self.TEXT_SAMPLE_BYTES]
        if not sample or b"\\x00" in sample:
            return False
        try:
            decoded = sample.decode("utf-8")
        except UnicodeDecodeError:
            try:
                decoded = sample.decode("latin-1")
            except Exception:
                return False
        if not decoded:
            return True
        printable = sum(
            ch.isprintable()
            or ch
            in "\\r\
\\t"
            for ch in decoded
        )
        return printable / len(decoded) >= 0.85

    def _extract_excel(self, file_path):
        logger.info(f"Reading Excel: {file_path}")
        from openpyxl import load_workbook

        workbook = load_workbook(filename=file_path, data_only=True, read_only=True)
        content = []
        try:
            for sheet in workbook.sheetnames:
                content.append(
                    f"\
Sheet: {sheet}\
"
                )
                worksheet = workbook[sheet]
                for row in worksheet.iter_rows(values_only=True):
                    values = [str(cell) for cell in row if cell is not None]
                    if values:
                        content.append(" | ".join(values))
        finally:
            workbook.close()
        return ExtractedContent(
            text="\
".join(content),
            content_type=ContentType.XLSX,
            file_path=str(file_path),
            metadata={},
        )

    def _extract_csv(self, file_path):
        logger.info(f"Reading CSV: {file_path}")
        return ExtractedContent(
            text=self._read_text(file_path),
            content_type=ContentType.CSV,
            file_path=str(file_path),
            metadata={},
        )

    def _extract_json(self, file_path):
        logger.info(f"Reading JSON: {file_path}")
        text = self._read_text(file_path)
        try:
            data = json.loads(text)
            normalized = json.dumps(data, indent=2, ensure_ascii=False)
            metadata = {"json_parsed": True}
            return ExtractedContent(
                text=normalized,
                content_type=ContentType.JSON,
                file_path=str(file_path),
                metadata=metadata,
            )
        except json.JSONDecodeError as exc:
            # Preserve useful configuration even when it is JSON-like rather than strict JSON.
            logger.warning(
                f"Non-strict JSON; preserving raw text for {file_path}: {exc}"
            )
            return ExtractedContent(
                text=text,
                content_type=ContentType.JSON,
                file_path=str(file_path),
                metadata={"json_parsed": False, "json_parse_error": str(exc)},
            )

    def _extract_unknown_text(self, file_path):
        logger.info(f"Reading text with unknown extension: {file_path}")
        text = self._read_text(file_path)
        stripped = text.lstrip()
        if stripped.startswith(("{", "[")):
            try:
                data = json.loads(text)
                return ExtractedContent(
                    text=json.dumps(data, indent=2, ensure_ascii=False),
                    content_type=ContentType.JSON,
                    file_path=str(file_path),
                    metadata={"json_parsed": True, "source_extension": ""},
                )
            except json.JSONDecodeError:
                pass
        return ExtractedContent(
            text=text,
            content_type=ContentType.TEXT,
            file_path=str(file_path),
            metadata={"source_extension": ""},
        )

    def _extract_ldif(self, file_path):
        logger.info(f"Reading LDIF: {file_path}")
        return ExtractedContent(
            text=self._read_text(file_path),
            content_type=ContentType.LDIF,
            file_path=str(file_path),
            metadata={},
        )

    @staticmethod
    def _read_text(file_path: Path) -> str:
        return file_path.read_text(encoding="utf-8", errors="ignore")

    def _extract_text(self, file_path):
        logger.info(f"Reading text: {file_path}")
        return ExtractedContent(
            text=self._read_text(file_path),
            content_type=ContentType.TEXT,
            file_path=str(file_path),
            metadata={},
        )

    def _extract_pdf(self, file_path):
        logger.info(f"Reading PDF: {file_path}")
        import fitz

        document = fitz.open(file_path)
        try:
            content = [page.get_text() for page in document if page.get_text()]
        finally:
            document.close()
        return ExtractedContent(
            text="\
".join(content),
            content_type=ContentType.PDF,
            file_path=str(file_path),
            metadata={},
        )

    def _extract_docx(self, file_path):
        logger.info(f"Reading DOCX: {file_path}")
        from docx import Document

        document = Document(file_path)
        content = [p.text for p in document.paragraphs if p.text]
        return ExtractedContent(
            text="\
".join(content),
            content_type=ContentType.DOCX,
            file_path=str(file_path),
            metadata={},
        )

    def _extract_pptx(self, file_path):
        logger.info(f"Reading PPTX: {file_path}")
        from pptx import Presentation

        presentation = Presentation(file_path)
        content = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content.append(shape.text)
        return ExtractedContent(
            text="\
".join(content),
            content_type=ContentType.PPTX,
            file_path=str(file_path),
            metadata={},
        )

    def _extract_image(self, file_path):
        logger.info(f"Reading image: {file_path}")
        from PIL import Image, ImageFile
        import pytesseract

        try:
            image = Image.open(file_path)
            try:
                image.load()
                text = pytesseract.image_to_string(image)
                width, height = image.size
                image_format = image.format
            finally:
                image.close()
        except Exception as first_exc:
            # Some Confluence image downloads are valid but slightly truncated.
            # Give Pillow one controlled retry with truncated-image support.
            previous = ImageFile.LOAD_TRUNCATED_IMAGES
            try:
                ImageFile.LOAD_TRUNCATED_IMAGES = True
                image = Image.open(file_path)
                try:
                    image.load()
                    text = pytesseract.image_to_string(image)
                    width, height = image.size
                    image_format = image.format
                finally:
                    image.close()
            except Exception:
                # A few files are mislabeled image attachments containing SVG/XML.
                # Preserve useful searchable content rather than blindly trusting
                # the filename extension.
                ImageFile.LOAD_TRUNCATED_IMAGES = previous
                textual = self._extract_xml_image_fallback(file_path)
                if textual is not None:
                    return textual
                raise first_exc
            finally:
                ImageFile.LOAD_TRUNCATED_IMAGES = previous

        return ExtractedContent(
            text=text,
            content_type=ContentType.IMAGE,
            file_path=str(file_path),
            metadata={
                "width": width,
                "height": height,
                "format": image_format,
            },
        )

    @staticmethod
    def _extract_xml_image_fallback(file_path: Path):
        """Extract text from an image file that actually contains SVG/XML."""
        try:
            sample = file_path.read_bytes()[:65536]
            decoded = sample.decode("utf-8-sig", errors="strict")
        except (OSError, UnicodeDecodeError):
            return None

        stripped = decoded.lstrip().lower()
        if not (
            stripped.startswith("<?xml")
            or stripped.startswith("<svg")
            or "<svg" in stripped[:4096]
        ):
            return None

        text = file_path.read_text(encoding="utf-8-sig", errors="ignore")
        return ExtractedContent(
            text=text,
            content_type=ContentType.IMAGE,
            file_path=str(file_path),
            metadata={"format": "svg", "fallback": "xml-image"},
        )

    def _extract_zip(self, file_path):
        logger.info(f"Reading ZIP: {file_path}")
        content = []
        total_bytes = 0
        with zipfile.ZipFile(file_path, "r") as archive:
            infos = [info for info in archive.infolist() if not info.is_dir()]
            self._validate_archive_limits(infos, "ZIP")
            for info in infos:
                total_bytes += info.file_size
                if total_bytes > self.MAX_ARCHIVE_BYTES:
                    raise ValueError(
                        f"ZIP expanded size exceeds {self.MAX_ARCHIVE_BYTES} bytes"
                    )
                data = archive.read(info)
                decoded = self._decode_archive_member(data)
                if decoded is not None:
                    content.append(
                        f"\
FILE: {info.filename}\
"
                    )
                    content.append(decoded)
        return ExtractedContent(
            text="\
".join(content),
            content_type=ContentType.ZIP,
            file_path=str(file_path),
            metadata={"member_count": len(infos)},
        )

    def _extract_rar(self, file_path):
        logger.info(f"Reading RAR: {file_path}")
        import rarfile

        content = []
        total_bytes = 0
        with rarfile.RarFile(file_path) as archive:
            infos = [info for info in archive.infolist() if not info.is_dir()]
            self._validate_archive_limits(infos, "RAR")
            for info in infos:
                total_bytes += int(getattr(info, "file_size", 0) or 0)
                if total_bytes > self.MAX_ARCHIVE_BYTES:
                    raise ValueError(
                        f"RAR expanded size exceeds {self.MAX_ARCHIVE_BYTES} bytes"
                    )
                data = archive.read(info)
                decoded = self._decode_archive_member(data)
                if decoded is not None:
                    content.append(
                        f"\
FILE: {info.filename}\
"
                    )
                    content.append(decoded)
        return ExtractedContent(
            text="\
".join(content),
            content_type=ContentType.RAR,
            file_path=str(file_path),
            metadata={"member_count": len(infos)},
        )

    def _validate_archive_limits(self, infos, kind: str) -> None:
        if len(infos) > self.MAX_ARCHIVE_FILES:
            raise ValueError(
                f"{kind} contains {len(infos)} files; limit is {self.MAX_ARCHIVE_FILES}"
            )

    @staticmethod
    def _decode_archive_member(data: bytes) -> str | None:
        if b"\\x00" in data[:4096]:
            return None
        try:
            return data.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return data.decode("latin-1")
            except UnicodeDecodeError:
                return None
